"""PPTX parser built on python-pptx.

Each slide is one :class:`ParsedPage`. Slide title (when present), all shape
text, and the speaker-notes pane are concatenated into the page text. Table
shapes are extracted into ``page.tables``.

python-pptx exposes shape attributes through duck-typed accessors that
pyright cannot follow without help, so each call site casts the relevant
attribute to its concrete return type. This keeps every value flowing into
:mod:`packages.core.models` strictly typed.
"""

from __future__ import annotations

import asyncio
import logging
import re
from io import BytesIO
from typing import cast

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from packages.core.models.source import (
    ParsedPage,
    ParsedSource,
    SourceMetadataExtracted,
)
from packages.workers.source.parsers.lang_detect import detect_language

logger = logging.getLogger(__name__)


class PPTXParser:
    """Parses a PPTX byte stream into a :class:`ParsedSource`."""

    async def parse(self, file_bytes: bytes, filename: str) -> ParsedSource:
        return await asyncio.to_thread(self._parse_sync, file_bytes, filename)

    def _parse_sync(self, file_bytes: bytes, filename: str) -> ParsedSource:
        errors: list[str] = []

        try:
            prs = Presentation(BytesIO(file_bytes))
        except Exception as exc:
            logger.warning("Failed to open PPTX %s: %s", filename, exc)
            return ParsedSource(
                filename=filename,
                file_type="pptx",
                file_size_bytes=len(file_bytes),
                pages=[],
                metadata=SourceMetadataExtracted(),
                full_text="",
                needs_ocr_pages=[],
                parse_errors=[f"Failed to open PPTX: {exc}"],
            )

        pages: list[ParsedPage] = []
        for slide_index, slide in enumerate(prs.slides):
            try:
                pages.append(self._parse_slide(slide, slide_index))
            except Exception as exc:
                msg = f"Slide {slide_index + 1} parse error: {exc}"
                errors.append(msg)
                logger.warning(msg)

        full_text = "\n\n".join(p.text for p in pages if p.text)
        word_count = len(full_text.split()) if full_text else 0
        metadata = self._build_metadata(prs, full_text, word_count, len(pages))

        return ParsedSource(
            filename=filename,
            file_type="pptx",
            file_size_bytes=len(file_bytes),
            pages=pages,
            metadata=metadata,
            full_text=full_text,
            needs_ocr_pages=[],
            parse_errors=errors,
        )

    @staticmethod
    def _parse_slide(slide: object, slide_index: int) -> ParsedPage:
        title = ""
        body_parts: list[str] = []
        headings: list[str] = []
        tables: list[list[list[str]]] = []

        shapes_container = getattr(slide, "shapes", None)
        try:
            title_shape: object | None = (
                getattr(shapes_container, "title", None) if shapes_container is not None else None
            )
        except (AttributeError, KeyError):
            title_shape = None
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title = cast("str", getattr(title_shape, "text", "") or "").strip()
            if title:
                headings.append(title)

        shapes_iter = cast("list[object]", list(shapes_container) if shapes_container else [])
        for shape in shapes_iter:
            if getattr(shape, "has_table", False):
                rows_out: list[list[str]] = []
                shape_table = getattr(shape, "table", None)
                if shape_table is None:
                    continue
                for row in cast("list[object]", list(getattr(shape_table, "rows", []))):
                    cells = cast("list[object]", list(getattr(row, "cells", [])))
                    rows_out.append(
                        [cast("str", getattr(cell, "text", "") or "").strip() for cell in cells]
                    )
                if rows_out:
                    tables.append(rows_out)
                continue
            if getattr(shape, "has_text_frame", False):
                text_frame = getattr(shape, "text_frame", None)
                txt = cast("str", getattr(text_frame, "text", "") or "").strip()
                if txt and txt != title:
                    body_parts.append(txt)

        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_slide = getattr(slide, "notes_slide", None)
                notes_text_frame = (
                    getattr(notes_slide, "notes_text_frame", None) if notes_slide else None
                )
                notes_text = cast("str", getattr(notes_text_frame, "text", "") or "").strip()
            except (AttributeError, KeyError):
                notes_text = ""

        sections: list[str] = []
        if title:
            sections.append(title)
        sections.extend(body_parts)
        if notes_text:
            sections.append(f"[Speaker notes]\n{notes_text}")
        page_text = "\n\n".join(sections)
        char_count = len(re.sub(r"\s+", "", page_text))

        return ParsedPage(
            page_number=slide_index + 1,
            text=page_text,
            char_count=char_count,
            needs_ocr=False,
            headings=headings,
            tables=tables,
        )

    @staticmethod
    def _build_metadata(
        prs: PresentationType,
        full_text: str,
        word_count: int,
        page_count: int,
    ) -> SourceMetadataExtracted:
        props = getattr(prs, "core_properties", None)
        title: str | None = None
        authors: list[str] = []
        creation_date: str | None = None
        year: int | None = None

        if props is not None:
            title = (props.title or "").strip() or None
            author_field = (props.author or "").strip()
            authors = [a.strip() for a in re.split(r"[,;]", author_field) if a.strip()]
            created = props.created
            creation_date = created.isoformat() if created else None
            year = created.year if created else None

        return SourceMetadataExtracted(
            title=title,
            authors=authors,
            year=year,
            doi=None,
            page_count=page_count,
            word_count=word_count,
            language_detected=detect_language(full_text),
            has_images=False,
            creation_date=creation_date,
        )
