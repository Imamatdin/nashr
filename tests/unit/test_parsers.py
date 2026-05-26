"""Unit tests for the per-format parsers and the routing service.

These run against the real libraries (PyMuPDF, python-docx, python-pptx,
Pillow). For test cases that require a PDF with specific properties (a
scanned-only page, a heading-vs-body font split), we generate them inline
with reportlab so tests/golden/ stays minimal.
"""

from __future__ import annotations

import io
import shutil
from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptxPt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import (
    Image as RLImage,
)
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

from packages.workers.source import SourceParseService
from packages.workers.source.parsers import (
    DOCXParser,
    ImageParser,
    PDFParser,
    PPTXParser,
    TextParser,
    XLSXParser,
    detect_language,
)

GOLDEN = Path(__file__).resolve().parent.parent / "golden"

TESSERACT_AVAILABLE = shutil.which("tesseract") is not None or any(
    Path(p).exists()
    for p in (
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    )
)


# ---------------------------------------------------------------------------
# language detection
# ---------------------------------------------------------------------------


def test_detect_uzbek_latin() -> None:
    text = (
        "Ag'artıwshılıq XVIII asirde payda boldı. Volter, Russo va Monteske "
        "oz oylari bilan jamiyatni ozgartirdi. Bu davrning oqishlari, "
        "g'oyalari va shuhrat-shoxratlari hozirgi kungacha yetib keldi."
    )
    assert detect_language(text) == "uz"


def test_detect_russian_cyrillic() -> None:
    text = (
        "Просвещение возникло в XVIII веке в Европе. Вольтер, Руссо и "
        "Монтескье изменили общество своими идеями о разуме и свободе."
    )
    assert detect_language(text) == "ru"


def test_detect_english() -> None:
    text = (
        "The Enlightenment was an intellectual movement in eighteenth century "
        "Europe. Thinkers like Voltaire and Rousseau transformed society with "
        "ideas about reason, liberty, and the social contract."
    )
    assert detect_language(text) == "en"


def test_detect_mixed_defaults_to_majority() -> None:
    text = "The Enlightenment Просвещение возникло в Европе и охватило весь континент."
    assert detect_language(text) == "ru"


def test_detect_empty_returns_none() -> None:
    assert detect_language("") is None
    assert detect_language("12345 !@#$%") is None


# ---------------------------------------------------------------------------
# PDF parser — fixture builders
# ---------------------------------------------------------------------------


def _make_pdf_with_image_page() -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, title="mixed pdf")
    styles = getSampleStyleSheet()
    story: list[object] = [
        Paragraph(
            "This first page contains a generous helping of textual content. "
            "It is long enough to clear the text-density threshold of fifty "
            "characters that the parser uses, so it should NOT be flagged as "
            "needing OCR by the heuristic. The second page is a pure image.",
            styles["BodyText"],
        ),
        PageBreak(),
    ]
    image_buffer = BytesIO()
    Image.new("RGB", (300, 300), color=(220, 30, 30)).save(image_buffer, format="PNG")
    image_buffer.seek(0)
    story.append(RLImage(image_buffer, width=200, height=200))
    doc.build(story)
    return buffer.getvalue()


def _make_pdf_with_captioned_figure() -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, title="figure pdf")
    styles = getSampleStyleSheet()
    image_buffer = BytesIO()
    Image.new("RGB", (300, 300), color=(20, 120, 160)).save(image_buffer, format="PNG")
    image_buffer.seek(0)
    story: list[object] = [
        Paragraph(
            "Datacenter thermal management is the subject of this section, with "
            "enough body text to comfortably clear the OCR density threshold used "
            "by the parser heuristic so the page is treated as real text.",
            styles["BodyText"],
        ),
        RLImage(image_buffer, width=200, height=200),
        Paragraph(
            "Figure 1: A supercritical CO2 cooling loop for a server rack.",
            styles["BodyText"],
        ),
    ]
    doc.build(story)
    return buffer.getvalue()


def _make_pdf_with_tiny_image() -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, title="tiny image pdf")
    styles = getSampleStyleSheet()
    image_buffer = BytesIO()
    Image.new("RGB", (40, 40), color=(0, 0, 0)).save(image_buffer, format="PNG")
    image_buffer.seek(0)
    story: list[object] = [
        Paragraph(
            "Body text long enough to parse cleanly and clear the density check "
            "while the only embedded image is a tiny decorative icon below the "
            "figure threshold the parser enforces.",
            styles["BodyText"],
        ),
        RLImage(image_buffer, width=20, height=20),
    ]
    doc.build(story)
    return buffer.getvalue()


def _make_pdf_with_headings() -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, title="heading pdf")
    styles = getSampleStyleSheet()
    heading = ParagraphStyle(
        "BigHeading",
        parent=styles["Heading1"],
        fontSize=18,
        leading=22,
        spaceAfter=8,
    )
    body = ParagraphStyle(
        "BodyTwelve",
        parent=styles["BodyText"],
        fontSize=12,
        leading=14,
    )
    story: list[object] = [
        Paragraph("Important Section Heading", heading),
        Spacer(1, 6),
        Paragraph(
            "This is body text rendered at twelve points. It must NOT be "
            "classified as a heading by the parser's font-size heuristic. "
            "The heading above is at eighteen points and bold.",
            body,
        ),
    ]
    doc.build(story)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# PDF parser
# ---------------------------------------------------------------------------


async def test_pdf_parse_extracts_text() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "sample_3page.pdf").read_bytes()

    result = await parser.parse(payload, "sample_3page.pdf")

    assert len(result.pages) == 3
    assert all(page.char_count > 100 for page in result.pages)
    # Use ASCII-safe tokens unique to this fixture: reportlab's default font
    # substitutes some Uzbek diacritics during rendering, so we assert against
    # content that survives font fallback.
    assert "Volter" in result.full_text
    assert "Frantsiya" in result.full_text
    assert result.metadata.page_count == 3
    assert result.metadata.word_count > 500


async def test_pdf_parse_detects_scanned_pages() -> None:
    parser = PDFParser()
    payload = _make_pdf_with_image_page()

    result = await parser.parse(payload, "mixed.pdf")

    assert len(result.pages) == 2
    assert result.pages[0].needs_ocr is False
    assert result.pages[1].needs_ocr is True
    assert 2 in result.needs_ocr_pages


async def test_pdf_parse_extracts_figure_with_caption() -> None:
    # Image engine source grounding: embedded raster + its "Figure N:" caption
    # + surrounding page text, so the engine can topic-match a slide subject.
    parser = PDFParser()
    result = await parser.parse(_make_pdf_with_captioned_figure(), "figure.pdf")

    assert len(result.figures) == 1
    figure = result.figures[0]
    assert figure.content_type.startswith("image/")
    assert len(figure.data) > 0
    assert figure.width >= 150
    assert figure.height >= 150
    assert figure.page_number == 1
    assert figure.caption is not None
    assert "supercritical CO2 cooling loop" in figure.caption
    # The broader page text rides along as context for topic matching.
    assert "thermal management" in figure.context


async def test_pdf_parse_skips_tiny_images() -> None:
    parser = PDFParser()
    result = await parser.parse(_make_pdf_with_tiny_image(), "tiny.pdf")
    assert result.figures == []


async def test_pdf_parse_no_figures_when_no_images() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "sample_3page.pdf").read_bytes()
    result = await parser.parse(payload, "sample_3page.pdf")
    assert result.figures == []


async def test_pdf_parse_extracts_doi() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "sample_with_doi.pdf").read_bytes()

    result = await parser.parse(payload, "sample_with_doi.pdf")

    assert result.metadata.doi is not None
    assert result.metadata.doi.startswith("10.")


async def test_pdf_parse_detects_uzbek_language() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "sample_3page.pdf").read_bytes()

    result = await parser.parse(payload, "sample_3page.pdf")

    assert result.metadata.language_detected == "uz"


async def test_pdf_parse_handles_empty_pdf() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "empty.pdf").read_bytes()

    result = await parser.parse(payload, "empty.pdf")

    assert len(result.pages) == 1
    assert result.pages[0].char_count == 0
    assert result.metadata.word_count == 0


async def test_pdf_parse_handles_prompt_injection_as_normal_text() -> None:
    parser = PDFParser()
    payload = (GOLDEN / "prompt_injection.pdf").read_bytes()

    result = await parser.parse(payload, "prompt_injection.pdf")

    assert "IGNORE ALL PREVIOUS INSTRUCTIONS" in result.full_text


async def test_pdf_parse_extracts_headings() -> None:
    parser = PDFParser()
    payload = _make_pdf_with_headings()

    result = await parser.parse(payload, "heading.pdf")

    assert len(result.pages) == 1
    headings = result.pages[0].headings
    joined = " ".join(headings)
    assert "Important Section Heading" in joined
    assert all("body text" not in h.lower() for h in headings)


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------


def _make_docx_with_table() -> bytes:
    document = Document()
    document.add_heading("Test document", level=1)
    document.add_paragraph("Body intro.")
    table = document.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "name"
    table.cell(0, 1).text = "year"
    table.cell(0, 2).text = "role"
    table.cell(1, 0).text = "Volter"
    table.cell(1, 1).text = "1694"
    table.cell(1, 2).text = "philosopher"

    if document.styles["Normal"].font.size is None:
        document.styles["Normal"].font.size = Pt(12)

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


async def test_docx_parse_extracts_text() -> None:
    parser = DOCXParser()
    payload = (GOLDEN / "sample_article.docx").read_bytes()

    result = await parser.parse(payload, "sample_article.docx")

    assert "Kirish" in result.full_text
    assert result.metadata.word_count > 100


async def test_docx_parse_groups_by_headings() -> None:
    parser = DOCXParser()
    payload = (GOLDEN / "sample_article.docx").read_bytes()

    result = await parser.parse(payload, "sample_article.docx")

    assert len(result.pages) >= 3
    assert any(page.headings for page in result.pages)


async def test_docx_parse_extracts_tables() -> None:
    parser = DOCXParser()
    payload = _make_docx_with_table()

    result = await parser.parse(payload, "with_table.docx")

    assert result.pages, "expected at least one page"
    table_pages = [p for p in result.pages if p.tables]
    assert table_pages, "expected at least one page to carry the table"
    table = table_pages[0].tables[0]
    assert table[0] == ["name", "year", "role"]
    assert table[1] == ["Volter", "1694", "philosopher"]


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


def _make_pptx(with_notes: bool = False) -> bytes:
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title only
    slides_text = ["First slide title", "Second slide title", "Third slide title"]
    for index, title_text in enumerate(slides_text):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title_text
        body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
        tf = body.text_frame
        tf.text = f"Body content for slide {index + 1}"
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxPt(18)
        if with_notes and index == 0:
            slide.notes_slide.notes_text_frame.text = "These are speaker notes for slide one."
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


async def test_pptx_parse_extracts_slides() -> None:
    parser = PPTXParser()
    payload = _make_pptx()

    result = await parser.parse(payload, "deck.pptx")

    assert len(result.pages) == 3
    for index, page in enumerate(result.pages, start=1):
        assert page.text, f"slide {index} produced empty text"
        assert f"slide {index}" in page.text.lower()


async def test_pptx_parse_extracts_speaker_notes() -> None:
    parser = PPTXParser()
    payload = _make_pptx(with_notes=True)

    result = await parser.parse(payload, "deck_with_notes.pptx")

    assert "speaker notes for slide one" in result.pages[0].text.lower()


# ---------------------------------------------------------------------------
# Text parser
# ---------------------------------------------------------------------------


async def test_text_parse_utf8() -> None:
    parser = TextParser()
    body = (
        "Ag'artıwshılıq dáwiri intellektual erkinlik hám sın oylawg'a "
        "tayanadı. Bul háreket Yevropanı ózgertti."
    )
    result = await parser.parse(body.encode("utf-8"), "note.txt")

    assert body in result.full_text
    assert result.metadata.language_detected == "uz"


async def test_text_parse_splits_long_text() -> None:
    parser = TextParser()
    body = "A" * 10_000

    result = await parser.parse(body.encode("utf-8"), "long.txt")

    assert len(result.pages) > 1
    assert sum(len(p.text) for p in result.pages) == 10_000


async def test_text_parse_latin1_fallback_does_not_raise() -> None:
    parser = TextParser()
    payload = b"Caf\xe9"  # invalid utf-8 byte for é

    result = await parser.parse(payload, "weird.txt")

    assert "Caf" in result.full_text
    assert any("latin-1" in err for err in result.parse_errors)


# ---------------------------------------------------------------------------
# XLSX parser
# ---------------------------------------------------------------------------


def _make_xlsx_two_sheets() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "People"
    sheet.append(["name", "born", "school"])
    sheet.append(["Volter", 1694, "Ag'artıwshılıq"])
    sheet.append(["Russo", 1712, "Ag'artıwshılıq"])

    works = workbook.create_sheet(title="Works")
    works.append(["title", "year"])
    works.append(["Kandid", 1759])

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _make_xlsx_empty() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Empty"
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


async def test_xlsx_parse_extracts_text() -> None:
    parser = XLSXParser()
    payload = _make_xlsx_two_sheets()

    result = await parser.parse(payload, "people.xlsx")

    assert len(result.pages) == 2
    assert result.pages[0].headings == ["People"]
    assert result.pages[1].headings == ["Works"]
    assert "Volter" in result.full_text
    assert "1694" in result.full_text
    assert "Kandid" in result.full_text
    assert result.metadata.word_count > 0
    assert result.needs_ocr_pages == []


async def test_xlsx_parse_extracts_tables() -> None:
    parser = XLSXParser()
    payload = _make_xlsx_two_sheets()

    result = await parser.parse(payload, "people.xlsx")

    table = result.pages[0].tables[0]
    assert table[0] == ["name", "born", "school"]
    assert table[1] == ["Volter", "1694", "Ag'artıwshılıq"]
    assert table[2] == ["Russo", "1712", "Ag'artıwshılıq"]


async def test_xlsx_parse_handles_empty_workbook() -> None:
    parser = XLSXParser()
    payload = _make_xlsx_empty()

    result = await parser.parse(payload, "empty.xlsx")

    assert result.parse_errors == []
    # A workbook with no data rows still produces one ParsedPage so downstream
    # code has a stable shape; the page carries no extracted tables.
    assert len(result.pages) == 1
    assert result.pages[0].tables == []
    # No body rows means no needs_ocr signal and no errors.
    assert result.needs_ocr_pages == []


async def test_xlsx_parse_handles_corrupt_bytes() -> None:
    parser = XLSXParser()

    result = await parser.parse(b"not a real zip archive", "broken.xlsx")

    assert result.parse_errors, "corrupt input must surface a parse error"
    assert result.full_text == ""


async def test_xlsx_parse_uses_golden_fixture() -> None:
    parser = XLSXParser()
    payload = (GOLDEN / "sample_spreadsheet.xlsx").read_bytes()

    result = await parser.parse(payload, "sample_spreadsheet.xlsx")

    assert len(result.pages) == 2
    assert "Volter" in result.full_text
    assert "Kandid" in result.full_text


# ---------------------------------------------------------------------------
# Image parser
# ---------------------------------------------------------------------------


async def test_image_parse_marks_ocr_needed() -> None:
    parser = ImageParser()
    payload = (GOLDEN / "sample_scanned.png").read_bytes()

    result = await parser.parse(payload, "sample_scanned.png")

    assert result.pages[0].needs_ocr is True
    assert result.pages[0].text == ""
    assert result.needs_ocr_pages == [1]


async def test_image_parse_records_dimensions_in_metadata() -> None:
    parser = ImageParser()
    image = Image.new("RGB", (640, 480), color=(0, 0, 0))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")

    result = await parser.parse(buffer.getvalue(), "blank.png")

    assert result.metadata.has_images is True
    assert result.metadata.title is not None
    assert "640x480" in result.metadata.title


# ---------------------------------------------------------------------------
# Parse-service routing
# ---------------------------------------------------------------------------


@pytest.fixture(scope="module")
def parse_service() -> SourceParseService:
    return SourceParseService()


async def test_parse_service_routes_pdf(parse_service: SourceParseService) -> None:
    payload = (GOLDEN / "sample_3page.pdf").read_bytes()
    result = await parse_service.parse(payload, "test.pdf", "pdf")
    assert result.file_type == "pdf"
    assert len(result.pages) == 3


async def test_parse_service_routes_docx(parse_service: SourceParseService) -> None:
    payload = (GOLDEN / "sample_article.docx").read_bytes()
    result = await parse_service.parse(payload, "test.docx", "docx")
    assert result.file_type == "docx"


async def test_parse_service_routes_pptx(parse_service: SourceParseService) -> None:
    result = await parse_service.parse(_make_pptx(), "test.pptx", "pptx")
    assert result.file_type == "pptx"
    assert len(result.pages) == 3


async def test_parse_service_routes_xlsx(parse_service: SourceParseService) -> None:
    payload = (GOLDEN / "sample_spreadsheet.xlsx").read_bytes()
    result = await parse_service.parse(payload, "test.xlsx", "xlsx")
    assert result.file_type == "xlsx"
    assert len(result.pages) == 2
    assert "Volter" in result.full_text


async def test_parse_service_routes_txt(parse_service: SourceParseService) -> None:
    result = await parse_service.parse(b"hello world", "n.txt", "txt")
    assert result.file_type == "txt"
    assert "hello world" in result.full_text


@pytest.mark.skipif(not TESSERACT_AVAILABLE, reason="Tesseract not installed")
async def test_parse_service_routes_image_with_ocr(
    parse_service: SourceParseService,
) -> None:
    """With Tesseract installed, the routing service runs OCR end-to-end."""

    payload = (GOLDEN / "sample_scanned.png").read_bytes()
    result = await parse_service.parse(payload, "scan.png", "png")

    assert result.file_type == "png"
    assert len(result.pages) == 1
    assert result.needs_ocr_pages == []
    assert result.pages[0].is_ocr is True
    assert result.pages[0].text != ""


@pytest.mark.skipif(TESSERACT_AVAILABLE, reason="Only meaningful without Tesseract")
async def test_parse_service_routes_image_without_ocr(
    parse_service: SourceParseService,
) -> None:
    """Without Tesseract, the OCR signal is left unconsumed for downstream handling."""

    payload = (GOLDEN / "sample_scanned.png").read_bytes()
    result = await parse_service.parse(payload, "scan.png", "png")

    assert result.file_type == "png"
    assert len(result.pages) == 1
    assert result.needs_ocr_pages == [1]
    assert result.pages[0].text == ""


async def test_parse_service_rejects_unknown_type(
    parse_service: SourceParseService,
) -> None:
    with pytest.raises(ValueError, match="No parser is registered"):
        await parse_service.parse(b"data", "file.xyz", "xyz")
